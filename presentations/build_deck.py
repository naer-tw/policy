"""Build the 30-minute pastor advocacy presentation deck."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy
from lxml import etree

# ===== Design tokens =====
FONT_TITLE = "微軟正黑體"
FONT_BODY = "微軟正黑體"
FONT_SCRIPTURE = "標楷體"

NAVY = RGBColor(0x0B, 0x2B, 0x4A)        # deep biblical blue
GOLD = RGBColor(0xC9, 0xA2, 0x47)        # subdued gold
CREAM = RGBColor(0xF7, 0xF3, 0xEA)       # parchment
INK = RGBColor(0x1A, 0x1A, 0x1A)         # near black
RED = RGBColor(0xB2, 0x3A, 0x3A)         # alert red
GRAY = RGBColor(0x6B, 0x6B, 0x6B)
LIGHT_GOLD = RGBColor(0xE8, 0xD9, 0xA8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SOFT_BG = RGBColor(0xFA, 0xF7, 0xF1)

# 16:9
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

BLANK = prs.slide_layouts[6]


def set_run(run, text, font_name=FONT_BODY, size=24, bold=False, color=INK):
    run.text = text
    f = run.font
    f.name = font_name
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = color
    # Force East-Asian font setting (eastAsia) so CJK glyphs render with the chosen font
    rPr = run._r.get_or_add_rPr()
    # remove existing eastAsia entries
    for tag in ("ea", "cs", "latin"):
        for el in rPr.findall(qn(f"a:{tag}")):
            rPr.remove(el)
    ea = etree.SubElement(rPr, qn("a:ea"))
    ea.set("typeface", font_name)
    latin = etree.SubElement(rPr, qn("a:latin"))
    latin.set("typeface", font_name)


def add_bg(slide, color=CREAM):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    return bg


def add_text(slide, x, y, w, h, lines,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, fill=None, line=None):
    """lines: list of dict {text, size, bold, color, font, align}"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    if fill is not None:
        tb.fill.solid()
        tb.fill.fore_color.rgb = fill
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ln.get("align", align)
        if "space_after" in ln:
            p.space_after = Pt(ln["space_after"])
        if "space_before" in ln:
            p.space_before = Pt(ln["space_before"])
        run = p.add_run()
        set_run(
            run,
            ln["text"],
            font_name=ln.get("font", FONT_BODY),
            size=ln["size"],
            bold=ln.get("bold", False),
            color=ln.get("color", INK),
        )
    return tb


def add_rect(slide, x, y, w, h, fill=NAVY, line=None):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    r.fill.solid()
    r.fill.fore_color.rgb = fill
    if line is None:
        r.line.fill.background()
    else:
        r.line.color.rgb = line
        r.line.width = Pt(1)
    r.shadow.inherit = False
    return r


def add_line(slide, x1, y1, x2, y2, color=GOLD, weight=2.5):
    ln = slide.shapes.add_connector(1, x1, y1, x2, y2)
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    return ln


def add_footer(slide, idx, total, section=""):
    add_text(
        slide,
        Inches(0.4), Inches(7.05), Inches(8), Inches(0.35),
        [{"text": f"國教行動聯盟  ·  跨出牆外的服事  ·  {section}",
          "size": 10, "color": GRAY}],
    )
    add_text(
        slide,
        Inches(11.5), Inches(7.05), Inches(1.5), Inches(0.35),
        [{"text": f"{idx} / {total}",
          "size": 10, "color": GRAY, "align": PP_ALIGN.RIGHT}],
        align=PP_ALIGN.RIGHT,
    )


def add_section_chip(slide, text, color=NAVY):
    chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(0.4), Inches(2.6), Inches(0.45))
    chip.adjustments[0] = 0.5
    chip.fill.solid()
    chip.fill.fore_color.rgb = color
    chip.line.fill.background()
    tf = chip.text_frame
    tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    set_run(run, text, size=14, bold=True, color=WHITE)


def add_notes(slide, text):
    notes = slide.notes_slide.notes_text_frame
    notes.text = ""
    for i, para in enumerate(text.split("\n")):
        p = notes.paragraphs[0] if i == 0 else notes.add_paragraph()
        r = p.add_run()
        set_run(r, para, size=12, color=INK)


TOTAL = 17

# =============================================================
# Slide 1 — Cover
# =============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, NAVY)
# Decorative side band
side = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.35), SH)
side.fill.solid(); side.fill.fore_color.rgb = GOLD; side.line.fill.background()

add_text(s, Inches(1.0), Inches(0.7), Inches(11), Inches(0.5),
         [{"text": "國教行動聯盟  Action Alliance on Basic Education",
           "size": 14, "color": LIGHT_GOLD, "bold": False}])

add_text(s, Inches(1.0), Inches(1.5), Inches(11.3), Inches(2.4),
         [{"text": "跨出牆外的服事", "size": 60, "bold": True, "color": WHITE}])

add_text(s, Inches(1.0), Inches(3.4), Inches(11.3), Inches(2.0),
         [{"text": "從家庭崩解的呼喊", "size": 32, "color": CREAM},
          {"text": "到公共政策的呼召", "size": 32, "color": CREAM}])

# Scripture strip
add_rect(s, Inches(1.0), Inches(5.4), Inches(11.3), Inches(1.0), fill=RGBColor(0x14, 0x3A, 0x60))
add_text(s, Inches(1.2), Inches(5.45), Inches(11.0), Inches(0.95),
         [{"text": "「約瑟……他的枝條探出牆外。」",
           "size": 22, "color": LIGHT_GOLD, "font": FONT_SCRIPTURE, "bold": True},
          {"text": "—— 創世記 49:22",
           "size": 14, "color": GOLD}])

add_text(s, Inches(1.0), Inches(6.7), Inches(11.3), Inches(0.4),
         [{"text": "牧者聚集  ·  30 分鐘簡報  ·  2026", "size": 12, "color": LIGHT_GOLD}])

add_notes(s, "【影片開場 — 不在此片】\n\n播放國教盟 2023-2025 倡議成果影片（2 分鐘）。\n\n影片結束→第一句話：\n『您剛剛看見的，是過去三年神在牆外動工的痕跡。但今天我要先帶各位回到兩個現場——一個叫剴剴，一個叫性平法。因為這兩個現場，正在告訴台灣教會：我們再也不能只在牆內服事了。』")

# =============================================================
# Slide 2 — 剴剴
# =============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, NAVY)
add_section_chip(s, "現場真實  ①", color=RED)

# Top number / label
add_text(s, Inches(0.6), Inches(1.2), Inches(12.2), Inches(0.8),
         [{"text": "剴剴", "size": 64, "bold": True, "color": WHITE,
           "align": PP_ALIGN.CENTER}], align=PP_ALIGN.CENTER)

add_text(s, Inches(0.6), Inches(2.4), Inches(12.2), Inches(0.6),
         [{"text": "那個沒有人來抱起他的孩子",
           "size": 24, "color": LIGHT_GOLD, "align": PP_ALIGN.CENTER}],
         align=PP_ALIGN.CENTER)

# decorative bar
add_line(s, Inches(5.5), Inches(3.4), Inches(7.83), Inches(3.4), color=GOLD, weight=3)

add_text(s, Inches(0.8), Inches(3.9), Inches(11.7), Inches(2.0),
         [{"text": "他只有 1 歲。",
           "size": 30, "color": WHITE, "align": PP_ALIGN.CENTER},
          {"text": "他活下來的機會，",
           "size": 30, "color": WHITE, "align": PP_ALIGN.CENTER},
          {"text": "本來叫做「社會安全網」。",
           "size": 30, "color": GOLD, "bold": True, "align": PP_ALIGN.CENTER}],
         align=PP_ALIGN.CENTER)

add_footer(s, 2, TOTAL, "現場真實")
add_notes(s, "【講稿｜90 秒】\n\n2023 年的冬天，台北一個 1 歲多的男孩，我們叫他剴剴。他原本應該在保母家被照顧。社工去訪視過，那本來是社會安全網最後一道防線。\n\n但剴剴沒有等到下一次訪視——他在那雙本該擁抱他的手裡，永遠閉上了眼睛。\n\n弟兄姊妹，整個台灣震驚的不只是一個孩子的死，而是震驚我們以為的『系統』，原來這麼薄。社工被起訴、一審判刑——這在台灣是史無前例的。\n\n但判完之後呢？下一個剴剴會不會還在發生？")

# =============================================================
# Slide 3 — 性平法
# =============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, CREAM)
add_section_chip(s, "現場真實  ②", color=RED)

add_text(s, Inches(0.6), Inches(1.2), Inches(12.2), Inches(0.9),
         [{"text": "性平法", "size": 60, "bold": True, "color": NAVY,
           "align": PP_ALIGN.CENTER}], align=PP_ALIGN.CENTER)

add_text(s, Inches(0.6), Inches(2.3), Inches(12.2), Inches(0.6),
         [{"text": "當父母不再被允許說話",
           "size": 24, "color": GRAY, "align": PP_ALIGN.CENTER}],
         align=PP_ALIGN.CENTER)

# Two quote boxes
left = add_rect(s, Inches(1.0), Inches(3.4), Inches(5.5), Inches(2.6), fill=WHITE)
left.line.color.rgb = LIGHT_GOLD; left.line.width = Pt(1.5)
add_text(s, Inches(1.2), Inches(3.6), Inches(5.1), Inches(0.5),
         [{"text": "家長想問：", "size": 16, "color": GRAY}])
add_text(s, Inches(1.2), Inches(4.1), Inches(5.1), Inches(1.8),
         [{"text": "「我的孩子",
           "size": 26, "bold": True, "color": NAVY},
          {"text": " 在學什麼？」",
           "size": 26, "bold": True, "color": NAVY}])

right = add_rect(s, Inches(6.8), Inches(3.4), Inches(5.5), Inches(2.6), fill=NAVY)
add_text(s, Inches(7.0), Inches(3.6), Inches(5.1), Inches(0.5),
         [{"text": "制度卻說：", "size": 16, "color": LIGHT_GOLD}])
add_text(s, Inches(7.0), Inches(4.1), Inches(5.1), Inches(1.8),
         [{"text": "「這不是",
           "size": 26, "bold": True, "color": WHITE},
          {"text": " 你的事。」",
           "size": 26, "bold": True, "color": GOLD}])

add_text(s, Inches(0.6), Inches(6.2), Inches(12.2), Inches(0.5),
         [{"text": "誰來定義我們的孩子，應該長成什麼樣子？",
           "size": 18, "color": RED, "bold": True, "align": PP_ALIGN.CENTER}],
         align=PP_ALIGN.CENTER)

add_footer(s, 3, TOTAL, "現場真實")
add_notes(s, "【講稿｜90 秒】\n\n另一個現場是性平法。各位牧者，這不只是某一個爭議議題，而是整個世代的價值權柄之爭——誰來定義我們的孩子應該長成什麼樣子？\n\n過去兩年，我們看見家長在學校門口拿到的，是他們從來沒見過的教材；想去公聽會發言，被排除在程序之外。\n\n今年國教盟召開了『性平法再校準』公聽會，三大主張就是要把教育拉回教育的本位——讓家長重新有座位、讓孩子重新被守護。")

# =============================================================
# Slide 4 — 兩道呼喊
# =============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, CREAM)
add_section_chip(s, "現場真實  ③", color=NAVY)

add_text(s, Inches(0.6), Inches(1.0), Inches(12.2), Inches(0.8),
         [{"text": "這不是兩個個案", "size": 40, "color": GRAY,
           "align": PP_ALIGN.CENTER}], align=PP_ALIGN.CENTER)
add_text(s, Inches(0.6), Inches(1.7), Inches(12.2), Inches(0.8),
         [{"text": "是兩道呼喊", "size": 48, "bold": True, "color": NAVY,
           "align": PP_ALIGN.CENTER}], align=PP_ALIGN.CENTER)

# Two columns
col_y, col_h = Inches(3.0), Inches(2.2)
add_rect(s, Inches(1.0), col_y, Inches(5.3), col_h, fill=NAVY)
add_text(s, Inches(1.2), Inches(3.15), Inches(4.9), Inches(0.5),
         [{"text": "剴剴", "size": 24, "bold": True, "color": GOLD}])
add_text(s, Inches(1.2), Inches(3.7), Inches(4.9), Inches(0.6),
         [{"text": "兒少保護崩落", "size": 26, "bold": True, "color": WHITE}])
add_text(s, Inches(1.2), Inches(4.4), Inches(4.9), Inches(0.6),
         [{"text": "→ 系統的牆破了", "size": 20, "color": LIGHT_GOLD}])

add_rect(s, Inches(7.0), col_y, Inches(5.3), col_h, fill=NAVY)
add_text(s, Inches(7.2), Inches(3.15), Inches(4.9), Inches(0.5),
         [{"text": "性平法", "size": 24, "bold": True, "color": GOLD}])
add_text(s, Inches(7.2), Inches(3.7), Inches(4.9), Inches(0.6),
         [{"text": "教育價值傾斜", "size": 26, "bold": True, "color": WHITE}])
add_text(s, Inches(7.2), Inches(4.4), Inches(4.9), Inches(0.6),
         [{"text": "→ 文化的牆斜了", "size": 20, "color": LIGHT_GOLD}])

# Big question
add_rect(s, Inches(1.0), Inches(5.7), Inches(11.3), Inches(0.95), fill=GOLD)
add_text(s, Inches(1.0), Inches(5.7), Inches(11.3), Inches(0.95),
         [{"text": "教會還能只在牆內服事嗎？",
           "size": 30, "bold": True, "color": NAVY, "align": PP_ALIGN.CENTER}],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_footer(s, 4, TOTAL, "現場真實")
add_notes(s, "【講稿｜60 秒】\n\n一個是『保護的牆破了』，一個是『價值的牆斜了』。這就是 2026 年台灣兒少與家庭的處境。\n\n而教會——我們是這個社會剩下最完整的一張網。但我們需要走出去。")

# =============================================================
# Slide 5 — 三個數字
# =============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, CREAM)
add_section_chip(s, "成果聚焦  ①", color=NAVY)

add_text(s, Inches(0.6), Inches(1.0), Inches(12.2), Inches(0.7),
         [{"text": "神已經動工了", "size": 38, "bold": True, "color": NAVY,
           "align": PP_ALIGN.CENTER}], align=PP_ALIGN.CENTER)
add_text(s, Inches(0.6), Inches(1.75), Inches(12.2), Inches(0.5),
         [{"text": "2023 - 2025  三個數字", "size": 18, "color": GRAY,
           "align": PP_ALIGN.CENTER}], align=PP_ALIGN.CENTER)

# Three number cards
nums = [("113", "篇新聞稿"), ("286", "場活動"), ("168+", "合作組織")]
x0 = Inches(0.9)
gap = Inches(0.3)
card_w = Inches(3.85)
card_h = Inches(3.0)
for i, (n, lbl) in enumerate(nums):
    x = Emu(int(x0) + i * (int(card_w) + int(gap)))
    add_rect(s, x, Inches(2.55), card_w, card_h, fill=NAVY)
    # Gold underline
    add_line(s, Emu(int(x) + Inches(0.6).emu), Inches(4.5),
             Emu(int(x) + int(card_w) - Inches(0.6).emu), Inches(4.5),
             color=GOLD, weight=3)
    add_text(s, x, Inches(2.85), card_w, Inches(1.6),
             [{"text": n, "size": 80, "bold": True, "color": WHITE,
               "align": PP_ALIGN.CENTER}], align=PP_ALIGN.CENTER)
    add_text(s, x, Inches(4.7), card_w, Inches(0.8),
             [{"text": lbl, "size": 22, "color": LIGHT_GOLD,
               "align": PP_ALIGN.CENTER}], align=PP_ALIGN.CENTER)

add_text(s, Inches(0.6), Inches(5.85), Inches(12.2), Inches(0.5),
         [{"text": "7,917+ 媒體聲量　·　校園安全 8 項訴求，6 項獲政府回應",
           "size": 16, "color": GRAY, "align": PP_ALIGN.CENTER}],
         align=PP_ALIGN.CENTER)

add_text(s, Inches(0.6), Inches(6.4), Inches(12.2), Inches(0.5),
         [{"text": "當有人願意跨出牆外，神就讓果子結出來。",
           "size": 18, "color": NAVY, "bold": True, "align": PP_ALIGN.CENTER}],
         align=PP_ALIGN.CENTER)

add_footer(s, 5, TOTAL, "成果聚焦")
add_notes(s, "【講稿｜60 秒】\n\n過去三年，神讓國教盟做的事用三個數字濃縮：113 篇新聞稿、286 場活動、168 個合作組織。校園安全 8 項訴求，6 項政府已經回應。\n\n這不是國教盟的功勞——這是當有人願意跨出牆外，神就讓果子結出來。")

# =============================================================
# Slide 6 — 尼希米的牆
# =============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, NAVY)
add_section_chip(s, "成果聚焦  ②", color=GOLD)

# Scripture centered
add_text(s, Inches(1.0), Inches(2.0), Inches(11.3), Inches(2.5),
         [{"text": "「這樣，我們修造城牆，",
           "size": 32, "color": WHITE, "font": FONT_SCRIPTURE,
           "align": PP_ALIGN.CENTER},
          {"text": "城牆就都連絡，",
           "size": 32, "color": WHITE, "font": FONT_SCRIPTURE,
           "align": PP_ALIGN.CENTER},
          {"text": "因為百姓齊心做工。」",
           "size": 36, "bold": True, "color": GOLD, "font": FONT_SCRIPTURE,
           "align": PP_ALIGN.CENTER}],
         align=PP_ALIGN.CENTER)

add_line(s, Inches(5.5), Inches(5.0), Inches(7.83), Inches(5.0), color=GOLD)

add_text(s, Inches(1.0), Inches(5.2), Inches(11.3), Inches(0.6),
         [{"text": "—— 尼希米記 4:6",
           "size": 18, "color": LIGHT_GOLD, "align": PP_ALIGN.CENTER}],
         align=PP_ALIGN.CENTER)

add_text(s, Inches(1.0), Inches(6.1), Inches(11.3), Inches(0.6),
         [{"text": "國教盟做的，是把「願意修牆的人」連起來。",
           "size": 18, "color": CREAM, "align": PP_ALIGN.CENTER}],
         align=PP_ALIGN.CENTER)

add_footer(s, 6, TOTAL, "成果聚焦")
add_notes(s, "【講稿｜60 秒】\n\n尼希米重建城牆的時候，沒有一個人能單獨修完整道牆。每一段都靠不同的家族、不同的職業、彼此緊鄰著做。\n\n國教盟這三年做的，就是這件事——把『願意修牆的人』連起來。今天我也邀請台灣的教會，加入這道城牆的修造。")

# =============================================================
# Slide 7 — 討論題 1
# =============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, GOLD)
add_section_chip(s, "小組討論  ①", color=NAVY)

add_text(s, Inches(0.6), Inches(1.0), Inches(12.2), Inches(0.7),
         [{"text": "先停下來，聽神在說什麼",
           "size": 36, "bold": True, "color": NAVY, "align": PP_ALIGN.CENTER}],
         align=PP_ALIGN.CENTER)
add_text(s, Inches(0.6), Inches(1.75), Inches(12.2), Inches(0.5),
         [{"text": "Discussion  ·  3 minutes", "size": 16, "color": NAVY,
           "align": PP_ALIGN.CENTER}], align=PP_ALIGN.CENTER)

# Big question card
add_rect(s, Inches(1.2), Inches(2.6), Inches(10.9), Inches(3.8), fill=WHITE)
add_text(s, Inches(1.5), Inches(2.9), Inches(10.3), Inches(3.4),
         [{"text": "剛剛影片裡的孩子、剴剴的案件、性平法的張力——",
           "size": 22, "color": INK, "space_after": 14},
          {"text": "哪一個畫面最觸動你？",
           "size": 30, "bold": True, "color": NAVY, "space_after": 18},
          {"text": "身為牧者，神正在如何透過這些「牆外」的呼喊，",
           "size": 22, "color": INK, "space_after": 6},
          {"text": "對你和你的會眾說話？",
           "size": 22, "color": INK}])

add_text(s, Inches(0.6), Inches(6.6), Inches(12.2), Inches(0.5),
         [{"text": "一個畫面 ＋ 一句感動  ·  每位約 30 秒  ·  共 3 分鐘",
           "size": 14, "color": NAVY, "align": PP_ALIGN.CENTER}],
         align=PP_ALIGN.CENTER)

add_footer(s, 7, TOTAL, "小組討論 1")
add_notes(s, "【主持引導｜約 20 秒】\n\n『不需要分析、不需要對策，只要分享：哪一個畫面神特別觸動我？計時 3 分鐘，每位 30 秒就夠。』\n\n計時器開始。")

# =============================================================
# Slide 8 — 兒童家庭部
# =============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, CREAM)
add_section_chip(s, "2026 重點  ①", color=NAVY)

add_text(s, Inches(0.6), Inches(1.0), Inches(12.2), Inches(0.7),
         [{"text": "2026 第一個戰場：兒童家庭部",
           "size": 32, "bold": True, "color": NAVY, "align": PP_ALIGN.CENTER}],
         align=PP_ALIGN.CENTER)

# Triangle of 3 stats
stats = [
    ("65%", "通報 5 年增加"),
    ("7%", "處置率僅"),
    ("こども", "對標日本\n家庭庁"),
]
positions = [(Inches(1.0), Inches(2.3)), (Inches(5.0), Inches(2.3)), (Inches(9.0), Inches(2.3))]
for (x, y), (n, lbl) in zip(positions, stats):
    add_rect(s, x, y, Inches(3.3), Inches(2.8), fill=NAVY)
    if n == "こども":
        add_text(s, x, y + Inches(0.4), Inches(3.3), Inches(1.4),
                 [{"text": n, "size": 48, "bold": True, "color": GOLD,
                   "align": PP_ALIGN.CENTER}], align=PP_ALIGN.CENTER)
    else:
        add_text(s, x, y + Inches(0.3), Inches(3.3), Inches(1.5),
                 [{"text": n, "size": 72, "bold": True, "color": GOLD,
                   "align": PP_ALIGN.CENTER}], align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(1.95), Inches(3.3), Inches(0.8),
             [{"text": lbl, "size": 16, "color": WHITE,
               "align": PP_ALIGN.CENTER}], align=PP_ALIGN.CENTER)

add_rect(s, Inches(1.0), Inches(5.45), Inches(11.3), Inches(1.0), fill=GOLD)
add_text(s, Inches(1.0), Inches(5.45), Inches(11.3), Inches(1.0),
         [{"text": "兒少權法第 7 條應升格二級「兒童家庭部」",
           "size": 24, "bold": True, "color": NAVY, "align": PP_ALIGN.CENTER}],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_footer(s, 8, TOTAL, "2026 重點 一")
add_notes(s, "【講稿｜75 秒】\n\n兒少通報案件五年內增加 65%，但實際處置率只有 7%——這代表 93% 的求救聲被吞掉了。\n\n日本去年成立了『こども家庭庁』，韓國也有。台灣呢？我們連一個專責部會都沒有。\n\n國教盟正在推動的，是讓兒少權法第 7 條那個只是『會報』的層級，升格為二級『兒童家庭部』——這是國家對下一代的態度宣告。")

# =============================================================
# Slide 9 — 兒少權法六大缺口
# =============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, CREAM)
add_section_chip(s, "2026 重點  ②", color=NAVY)

add_text(s, Inches(0.6), Inches(1.0), Inches(12.2), Inches(0.7),
         [{"text": "2026 第二個戰場：兒少權法六大缺口",
           "size": 30, "bold": True, "color": NAVY, "align": PP_ALIGN.CENTER}],
         align=PP_ALIGN.CENTER)

add_text(s, Inches(0.6), Inches(1.75), Inches(12.2), Inches(0.5),
         [{"text": "165 條草案  ×  六大制度盲點", "size": 18, "color": GRAY,
           "align": PP_ALIGN.CENTER}], align=PP_ALIGN.CENTER)

# 6 pills in 2 rows
labels = ["通報認定", "社工保證人地位", "安置責任",
          "兒少代表權", "資訊揭露", "跨部會權責"]
pill_w, pill_h = Inches(3.6), Inches(1.1)
gap_x, gap_y = Inches(0.3), Inches(0.3)
start_x = Inches(1.05)
start_y = Inches(2.7)
for i, lbl in enumerate(labels):
    row, col = divmod(i, 3)
    x = Emu(int(start_x) + col * (int(pill_w) + int(gap_x)))
    y = Emu(int(start_y) + row * (int(pill_h) + int(gap_y)))
    pill = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, pill_w, pill_h)
    pill.adjustments[0] = 0.4
    pill.fill.solid(); pill.fill.fore_color.rgb = NAVY
    pill.line.fill.background()
    # highlight #2 (social worker) — most critical
    if i == 1:
        pill.fill.fore_color.rgb = GOLD
        txt_color = NAVY
    else:
        txt_color = WHITE
    tf = pill.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    set_run(run, lbl, size=22, bold=True, color=txt_color)

add_text(s, Inches(0.6), Inches(5.85), Inches(12.2), Inches(0.5),
         [{"text": "剴剴案之後，整個社工體系在崩潰——",
           "size": 16, "color": GRAY, "align": PP_ALIGN.CENTER}],
         align=PP_ALIGN.CENTER)
add_text(s, Inches(0.6), Inches(6.3), Inches(12.2), Inches(0.5),
         [{"text": "因為法律對她們的責任邊界沒有寫清楚。",
           "size": 18, "color": NAVY, "bold": True, "align": PP_ALIGN.CENTER}],
         align=PP_ALIGN.CENTER)

add_footer(s, 9, TOTAL, "2026 重點 一")
add_notes(s, "【講稿｜75 秒】\n\n衛福部預告的兒少權法草案有 165 條，看起來很完整，但國教盟比對之後找出六大缺口。\n\n最關鍵的就是『社工保證人地位』——剴剴案的社工被判刑之後，整個社工體系在崩潰，因為法律對她們的責任邊界沒有寫清楚。\n\n我們不是要保護任何一方，我們是要把制度的牆重新砌實。")

# =============================================================
# Slide 10 — 為不能自辯者
# =============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, NAVY)
add_section_chip(s, "2026 重點  ③", color=GOLD)

add_text(s, Inches(1.0), Inches(2.3), Inches(11.3), Inches(2.5),
         [{"text": "「你當為不能自辯的開口，",
           "size": 36, "color": WHITE, "font": FONT_SCRIPTURE,
           "align": PP_ALIGN.CENTER, "space_after": 12},
          {"text": "為一切孤獨的伸冤。」",
           "size": 38, "bold": True, "color": GOLD, "font": FONT_SCRIPTURE,
           "align": PP_ALIGN.CENTER}],
         align=PP_ALIGN.CENTER)

add_line(s, Inches(5.5), Inches(5.0), Inches(7.83), Inches(5.0), color=GOLD)

add_text(s, Inches(1.0), Inches(5.2), Inches(11.3), Inches(0.6),
         [{"text": "—— 箴言 31:8",
           "size": 18, "color": LIGHT_GOLD, "align": PP_ALIGN.CENTER}],
         align=PP_ALIGN.CENTER)

add_text(s, Inches(1.0), Inches(6.1), Inches(11.3), Inches(0.6),
         [{"text": "兒少在公共政策桌上，永遠是不能自辯的那一群。",
           "size": 18, "color": CREAM, "align": PP_ALIGN.CENTER}],
         align=PP_ALIGN.CENTER)

add_footer(s, 10, TOTAL, "2026 重點 一")
add_notes(s, "【講稿｜30 秒】\n\n兒少在公共政策的桌上，永遠是『不能自辯』的那一群。今天教會若不為他們開口，誰開口？")

# =============================================================
# Slide 11 — 家庭韌性聯盟
# =============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, CREAM)
add_section_chip(s, "2026 重點  ④", color=NAVY)

add_text(s, Inches(0.6), Inches(1.0), Inches(12.2), Inches(0.7),
         [{"text": "橫向：家庭韌性行動聯盟",
           "size": 32, "bold": True, "color": NAVY, "align": PP_ALIGN.CENTER}],
         align=PP_ALIGN.CENTER)
add_text(s, Inches(0.6), Inches(1.75), Inches(12.2), Inches(0.5),
         [{"text": "從單一議題  →  家庭總體戰略",
           "size": 20, "color": GRAY, "align": PP_ALIGN.CENTER}],
         align=PP_ALIGN.CENTER)

# 6 hex labels in row
fields = ["婚姻", "生育", "教養", "照顧", "經濟", "信仰"]
fw = Inches(1.7)
fh = Inches(1.7)
total_w = len(fields) * int(fw) + (len(fields)-1) * int(Inches(0.2))
start = Emu((int(SW) - total_w) // 2)
for i, fld in enumerate(fields):
    x = Emu(int(start) + i * (int(fw) + int(Inches(0.2))))
    shp = s.shapes.add_shape(MSO_SHAPE.HEXAGON, x, Inches(2.7), fw, fh)
    shp.fill.solid(); shp.fill.fore_color.rgb = NAVY
    shp.line.color.rgb = GOLD; shp.line.width = Pt(2)
    tf = shp.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); set_run(run, fld, size=22, bold=True, color=WHITE)

add_rect(s, Inches(1.0), Inches(5.4), Inches(11.3), Inches(1.4), fill=GOLD)
add_text(s, Inches(1.2), Inches(5.5), Inches(10.9), Inches(1.2),
         [{"text": "「家庭不會因為議題分散就崩解，",
           "size": 22, "color": NAVY, "align": PP_ALIGN.CENTER, "space_after": 4},
          {"text": "它是被『整個生活擠垮的』。」",
           "size": 24, "bold": True, "color": NAVY, "align": PP_ALIGN.CENTER}],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_footer(s, 11, TOTAL, "2026 重點 二")
add_notes(s, "【講稿｜60 秒】\n\n過去倡議的盲點是『議題分散』——婚家、生育、教養、托育各自為政。\n\n家庭韌性行動聯盟要做的是把這些縫合起來，因為家庭從來不是被『某個政策』擊垮的，是被『整個社會的擠壓』擊垮的。")

# =============================================================
# Slide 12 — 國際研討會
# =============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, CREAM)
add_section_chip(s, "2026 重點  ⑤", color=NAVY)

add_text(s, Inches(0.6), Inches(1.0), Inches(12.2), Inches(0.7),
         [{"text": "縱向：藥物濫用防制國際研討會",
           "size": 30, "bold": True, "color": NAVY, "align": PP_ALIGN.CENTER}],
         align=PP_ALIGN.CENTER)
add_text(s, Inches(0.6), Inches(1.75), Inches(12.2), Inches(0.5),
         [{"text": "校園新興毒品  ×  國際對話",
           "size": 20, "color": GRAY, "align": PP_ALIGN.CENTER}],
         align=PP_ALIGN.CENTER)

# Three cards
items = [
    ("依托咪酯", "校園新興毒品", RED),
    ("喪屍煙彈", "化學戰等級攻擊", RED),
    ("唾液快篩\n二層檢測", "國際對話桌", GOLD),
]
x0 = Inches(0.9); card_w = Inches(3.85); gap = Inches(0.3)
for i, (n, lbl, accent) in enumerate(items):
    x = Emu(int(x0) + i * (int(card_w) + int(gap)))
    add_rect(s, x, Inches(2.7), card_w, Inches(3.0), fill=NAVY)
    # Top accent bar
    add_rect(s, x, Inches(2.7), card_w, Inches(0.18), fill=accent)
    add_text(s, x, Inches(3.0), card_w, Inches(1.7),
             [{"text": n, "size": 32, "bold": True, "color": WHITE,
               "align": PP_ALIGN.CENTER}],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x, Inches(4.8), card_w, Inches(0.7),
             [{"text": lbl, "size": 16, "color": LIGHT_GOLD,
               "align": PP_ALIGN.CENTER}], align=PP_ALIGN.CENTER)

add_text(s, Inches(0.6), Inches(6.2), Inches(12.2), Inches(0.5),
         [{"text": "這不是技術問題，這是為這個世代築牆。",
           "size": 22, "color": NAVY, "bold": True, "align": PP_ALIGN.CENTER}],
         align=PP_ALIGN.CENTER)

add_footer(s, 12, TOTAL, "2026 重點 二")
add_notes(s, "【講稿｜60 秒】\n\n我們的孩子現在面對的不是 30 年前的毒品。依托咪酯、喪屍煙彈、新型尼古丁袋——這些都是化學戰等級的攻擊。\n\n今年國教盟參與藥物濫用防制國際研討會，把『唾液快篩二層檢測機制』帶上國際對話桌。\n\n這不是技術問題，這是為這個世代築牆。")

# =============================================================
# Slide 13 — 公共政策小組（四象限）
# =============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, CREAM)
add_section_chip(s, "呼籲教會  ①", color=NAVY)

add_text(s, Inches(0.6), Inches(0.95), Inches(12.2), Inches(0.7),
         [{"text": "邀請：每一間教會的「公共政策小組」",
           "size": 28, "bold": True, "color": NAVY, "align": PP_ALIGN.CENTER}],
         align=PP_ALIGN.CENTER)
add_text(s, Inches(0.6), Inches(1.6), Inches(12.2), Inches(0.5),
         [{"text": "沒有教會需要做全部，但每一間教會都可以做一塊",
           "size": 16, "color": GRAY, "align": PP_ALIGN.CENTER}],
         align=PP_ALIGN.CENTER)

# 2x2 grid
quads = [
    ("代禱", "為立法者、社工、教師守望", "🙏"),
    ("研究", "聖經倫理 × 政策分析", "📚"),
    ("發聲", "公聽會、媒體、聯署", "📣"),
    ("陪伴", "受害家庭、第一線實務工作者", "🤝"),
]
qw = Inches(5.4); qh = Inches(2.0)
gx = Inches(0.4); gy = Inches(0.25)
start_x = Inches(1.0); start_y = Inches(2.35)
for i, (title, sub, _) in enumerate(quads):
    row, col = divmod(i, 2)
    x = Emu(int(start_x) + col * (int(qw) + int(gx)))
    y = Emu(int(start_y) + row * (int(qh) + int(gy)))
    add_rect(s, x, y, qw, qh, fill=NAVY)
    add_line(s,
             Emu(int(x) + Inches(0.3).emu), Emu(int(y) + Inches(1.0).emu),
             Emu(int(x) + Inches(0.7).emu), Emu(int(y) + Inches(1.0).emu),
             color=GOLD, weight=3)
    add_text(s, x, Emu(int(y) + Inches(0.2).emu), qw, Inches(0.7),
             [{"text": title, "size": 32, "bold": True, "color": GOLD}])
    add_text(s,
             Emu(int(x) + Inches(2.0).emu),
             Emu(int(y) + Inches(0.4).emu),
             Inches(3.2), Inches(1.4),
             [{"text": sub, "size": 16, "color": WHITE}],
             anchor=MSO_ANCHOR.MIDDLE)

add_text(s, Inches(0.6), Inches(6.85), Inches(12.2), Inches(0.4),
         [{"text": "找一塊您的教會最能委身的——這就是第一步",
           "size": 14, "color": NAVY, "bold": True, "align": PP_ALIGN.CENTER}],
         align=PP_ALIGN.CENTER)

add_footer(s, 13, TOTAL, "呼籲教會")
add_notes(s, "【講稿｜90 秒】\n\n我知道每一位牧者都很忙，所以我不是邀請各位『另起爐灶』——我是邀請各位在你的教會體質裡，找一塊你最能委身的。\n\n有些教會有禱告勇士，那就成為代禱那一塊；\n有些教會有律師、學者，那就成為研究那一塊；\n有些教會有發言力、媒體連結，那就成為發聲那一塊；\n有些教會有社工、輔導、家庭事工，那就成為陪伴那一塊。\n\n沒有任何一間教會需要做全部，但每一間教會都可以做一塊。")

# =============================================================
# Slide 14 — 現今的機會（以斯帖）
# =============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, NAVY)
add_section_chip(s, "呼籲教會  ②", color=GOLD)

add_text(s, Inches(1.0), Inches(2.0), Inches(11.3), Inches(2.8),
         [{"text": "「焉知你得了王后的位分，",
           "size": 34, "color": WHITE, "font": FONT_SCRIPTURE,
           "align": PP_ALIGN.CENTER, "space_after": 14},
          {"text": "不是為現今的機會嗎？」",
           "size": 38, "bold": True, "color": GOLD, "font": FONT_SCRIPTURE,
           "align": PP_ALIGN.CENTER}],
         align=PP_ALIGN.CENTER)

add_line(s, Inches(5.5), Inches(5.1), Inches(7.83), Inches(5.1), color=GOLD)

add_text(s, Inches(1.0), Inches(5.3), Inches(11.3), Inches(0.6),
         [{"text": "—— 以斯帖記 4:14",
           "size": 18, "color": LIGHT_GOLD, "align": PP_ALIGN.CENTER}],
         align=PP_ALIGN.CENTER)

add_text(s, Inches(1.0), Inches(6.1), Inches(11.3), Inches(0.6),
         [{"text": "教會被神放在這個時代——不是為了安全，是為了現今的機會。",
           "size": 18, "color": CREAM, "align": PP_ALIGN.CENTER}],
         align=PP_ALIGN.CENTER)

add_footer(s, 14, TOTAL, "呼籲教會")
add_notes(s, "【講稿｜45 秒】\n\n以斯帖如果留在王宮裡安全地當她的王后，整個猶太民族就沒了。她的『位分』，是為了『跨出去』。\n\n弟兄姊妹，台灣教會被神放在這個時代、這片土地，不是為了我們的安全，是為了現今的機會。")

# =============================================================
# Slide 15 — 約瑟金句
# =============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, NAVY)
# Side gold band
side = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.35), SH)
side.fill.solid(); side.fill.fore_color.rgb = GOLD; side.line.fill.background()

add_section_chip(s, "呼籲教會  ③", color=GOLD)

add_text(s, Inches(1.0), Inches(1.6), Inches(11.3), Inches(2.0),
         [{"text": "「約瑟……",
           "size": 50, "color": WHITE, "font": FONT_SCRIPTURE,
           "align": PP_ALIGN.CENTER},
          {"text": "他的枝條探出牆外。」",
           "size": 56, "bold": True, "color": GOLD, "font": FONT_SCRIPTURE,
           "align": PP_ALIGN.CENTER}],
         align=PP_ALIGN.CENTER)

add_text(s, Inches(1.0), Inches(4.0), Inches(11.3), Inches(0.5),
         [{"text": "—— 創世記 49:22",
           "size": 18, "color": LIGHT_GOLD, "align": PP_ALIGN.CENTER}],
         align=PP_ALIGN.CENTER)

# Gold divider
add_line(s, Inches(3.5), Inches(4.85), Inches(9.83), Inches(4.85), color=GOLD, weight=2)

add_text(s, Inches(1.0), Inches(5.1), Inches(11.3), Inches(2.0),
         [{"text": "教會的根，留在泉旁；",
           "size": 30, "bold": True, "color": WHITE,
           "align": PP_ALIGN.CENTER, "space_after": 12},
          {"text": "教會的果子，必須探出牆外。",
           "size": 30, "bold": True, "color": GOLD,
           "align": PP_ALIGN.CENTER}],
         align=PP_ALIGN.CENTER)

add_footer(s, 15, TOTAL, "呼籲教會")
add_notes(s, "【講稿｜45 秒｜全場高點，請放慢語速】\n\n弟兄姊妹，創世記四十九章二十二節說：『約瑟，他的枝條探出牆外。』\n\n（停一秒）\n\n教會的根，留在泉旁——我們不離開基督。\n但教會的果子，必須探出牆外。\n\n這就是國教盟的呼召，也是今天我們對每一間教會的邀請。")

# =============================================================
# Slide 16 — 討論題 2
# =============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, GOLD)
add_section_chip(s, "小組討論  ②", color=NAVY)

add_text(s, Inches(0.6), Inches(0.95), Inches(12.2), Inches(0.7),
         [{"text": "這一步，我們從哪裡開始？",
           "size": 36, "bold": True, "color": NAVY, "align": PP_ALIGN.CENTER}],
         align=PP_ALIGN.CENTER)
add_text(s, Inches(0.6), Inches(1.65), Inches(12.2), Inches(0.5),
         [{"text": "Discussion  ·  3 minutes", "size": 16, "color": NAVY,
           "align": PP_ALIGN.CENTER}], align=PP_ALIGN.CENTER)

# Question card
add_rect(s, Inches(1.0), Inches(2.4), Inches(11.3), Inches(4.0), fill=WHITE)

add_text(s, Inches(1.3), Inches(2.6), Inches(10.7), Inches(0.7),
         [{"text": "如果您的教會要踏出第一步參與公共政策——",
           "size": 22, "color": INK}])

# Three numbered lines
nums_q = [
    ("1", "是先從「代禱／研究／發聲／陪伴」哪一塊開始？"),
    ("2", "最大的攔阻是什麼？"),
    ("3", "您願意帶  誰  一起？"),
]
start_y = Inches(3.5)
for i, (n, q) in enumerate(nums_q):
    y = Emu(int(start_y) + i * Inches(0.85).emu)
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.6), y, Inches(0.6), Inches(0.6))
    circ.fill.solid(); circ.fill.fore_color.rgb = NAVY
    circ.line.fill.background()
    tf = circ.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); set_run(run, n, size=20, bold=True, color=GOLD)
    is_third = (i == 2)
    add_text(s, Inches(2.4), Emu(int(y) + Inches(0.05).emu),
             Inches(9.6), Inches(0.7),
             [{"text": q, "size": 22,
               "bold": is_third, "color": NAVY if is_third else INK}])

add_text(s, Inches(0.6), Inches(6.7), Inches(12.2), Inches(0.4),
         [{"text": "挑一題回應就好  ·  讓你想到的那個人，等一下你會打給他",
           "size": 14, "color": NAVY, "bold": True, "align": PP_ALIGN.CENTER}],
         align=PP_ALIGN.CENTER)

add_footer(s, 16, TOTAL, "小組討論 2")
add_notes(s, "【主持引導｜約 20 秒】\n\n『這次不需要全部回答，每位牧者挑一題回應就好。重點不是答得完整，是讓你想到的那個人，等一下你會打給他。』\n\n計時器開始。")

# =============================================================
# Slide 17 — 結語禱告
# =============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, NAVY)
side = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.35), SH)
side.fill.solid(); side.fill.fore_color.rgb = GOLD; side.line.fill.background()

add_section_chip(s, "結語禱告", color=GOLD)

add_text(s, Inches(1.0), Inches(1.4), Inches(11.3), Inches(0.8),
         [{"text": "為「探出牆外」的呼召禱告",
           "size": 30, "bold": True, "color": WHITE, "align": PP_ALIGN.CENTER}],
         align=PP_ALIGN.CENTER)

# Prayer card
add_rect(s, Inches(1.5), Inches(2.6), Inches(10.3), Inches(3.4),
         fill=RGBColor(0x14, 0x3A, 0x60))

add_text(s, Inches(1.8), Inches(2.85), Inches(9.7), Inches(3.0),
         [{"text": "主，求祢讓台灣教會的枝條，",
           "size": 24, "color": WHITE, "font": FONT_SCRIPTURE,
           "align": PP_ALIGN.CENTER, "space_after": 8},
          {"text": "長出牆外，",
           "size": 26, "bold": True, "color": GOLD, "font": FONT_SCRIPTURE,
           "align": PP_ALIGN.CENTER, "space_after": 12},
          {"text": "結出剴剴等不到的果子，",
           "size": 24, "color": WHITE, "font": FONT_SCRIPTURE,
           "align": PP_ALIGN.CENTER, "space_after": 8},
          {"text": "結出下一代需要的盼望。",
           "size": 24, "color": WHITE, "font": FONT_SCRIPTURE,
           "align": PP_ALIGN.CENTER}],
         align=PP_ALIGN.CENTER)

add_text(s, Inches(1.0), Inches(6.3), Inches(11.3), Inches(0.5),
         [{"text": "奉主耶穌的名禱告，阿們。",
           "size": 18, "color": LIGHT_GOLD, "align": PP_ALIGN.CENTER}],
         align=PP_ALIGN.CENTER)

add_footer(s, 17, TOTAL, "結語")
add_notes(s, "【結語禱告｜約 60 秒，邀請會眾起立或同心】\n\n主啊，謝謝祢讓我們聽見牆外的呼喊。\n\n求祢給台灣每一位牧者，今天回去之後，至少打一通電話、找到一位同行的人。\n\n讓祢的教會像約瑟那棵枝子——根扎在祢裡面，果子結在這片土地上。\n\n奉主耶穌的名禱告，阿們。")

# =============================================================
out = "/home/user/policy/presentations/跨出牆外的服事_30分鐘簡報.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Total slides: {len(prs.slides)}")
